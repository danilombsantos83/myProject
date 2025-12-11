from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor

# === Criar apresentação ===
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])  # Layout em branco

# === Função para adicionar retângulos de ação ===
def adicionar_acao(slide, left, top, width, height, texto):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(173, 216, 230)  # Azul claro
    shape.line.color.rgb = RGBColor(0, 0, 0)
    shape.text = texto
    shape.text_frame.paragraphs[0].font.size = Pt(12)
    return shape

# === Função para adicionar losango de decisão ===
def adicionar_decisao(slide, left, top, width, height, texto):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 153)  # Amarelo claro
    shape.line.color.rgb = RGBColor(0, 0, 0)
    shape.text = texto
    shape.text_frame.paragraphs[0].font.size = Pt(12)
    return shape

# === Função para conectar dois shapes com seta ===
def conectar(shapes_slide, shape_from, shape_to):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        shape_from.left + shape_from.width/2,
        shape_from.top + shape_from.height,
        shape_to.left + shape_to.width/2,
        shape_to.top
    )
    connector.line.color.rgb = RGBColor(0, 0, 0)
    return connector

# === Definir posições ===
top = 0.5
left = 1
width = 2
height = 0.8
espaco = 1.2

# === Adicionar caixas de ação e decisão ===
menu = adicionar_acao(slide, left, top, width, height, "Menu Principal\nEscolher ação")
top += espaco

atualizar = adicionar_acao(slide, left, top, width, height, "Atualizar Tabelas\n(Banco Binance → SQLite)")
limpar_backups = adicionar_acao(slide, left + 3, top, width, height, "Limpar Tabelas de Backup")
top += espaco

analise = adicionar_acao(slide, left, top, width, height, "Análise Técnica\n(Escolher par e intervalo, Backtest)")
exportar_excel = adicionar_acao(slide, left + 3, top, width, height, "Exportar Excel")
top += espaco

adicionar_par = adicionar_acao(slide, left, top, width, height, "Adicionar Novo Par de Moedas")
config = adicionar_acao(slide, left + 3, top, width, height, "Configurar Aplicativo\n(Caminho de Exports)")
top += espaco

sair = adicionar_acao(slide, left + 1.5, top, width, height, "Sair")

# === Conectar shapes ===
conectar(slide, menu, atualizar)
conectar(slide, menu, limpar_backups)
conectar(slide, menu, analise)
conectar(slide, menu, exportar_excel)
conectar(slide, menu, adicionar_par)
conectar(slide, menu, config)
conectar(slide, menu, sair)

# === Salvar apresentação ===
prs.save("Fluxograma_Aplicacao.pptx")
print("✅ Fluxograma gerado: Fluxograma_Aplicacao.pptx")
